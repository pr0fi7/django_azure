from django.shortcuts import render, redirect
from django.http import JsonResponse
import openai
import json

from django.contrib import auth
from django.contrib.auth.models import User
import openai.cli
from .models import Chat
from django.utils import timezone
import shutil

from .forms import DocumentUploadForm 
from django.http import JsonResponse    

from io import BytesIO
import pdfplumber
from langchain_openai import OpenAIEmbeddings
from langchain.vectorstores.chroma import Chroma
import os
from langchain_text_splitters import CharacterTextSplitter
from langchain.schema.document import Document
from hashlib import md5
import requests
from bs4 import BeautifulSoup

from django.contrib.auth.decorators import login_required

import docx
from pptx import Presentation
import zipfile
import pandas as pd
import base64
import cv2 
import numpy as np
from PIL import Image
from PyPDF2 import PdfReader

from dotenv import load_dotenv, find_dotenv
_ = load_dotenv(find_dotenv())
openai.api_key=os.getenv("OPENAI_API_KEY")
openaiapi = os.getenv("OPENAI_API_KEY")

CHROMA_PATH = "chroma.db"

def get_xls_text(xls_files):
    text = ""
    for xls_file in xls_files:
        xls = pd.ExcelFile(xls_file)
        for sheet in xls.sheet_names:
            df = pd.read_excel(xls, sheet)
            text += df.to_string()
    return text

def unzip_file(zip_file, extract_to):
    with zipfile.ZipFile(zip_file, 'r') as zip_ref:
        zip_ref.extractall(extract_to)


def load_image(image_file):
    # Ensure the image file is opened correctly
    image = Image.open(image_file)
    image = image.convert('RGB')  # Ensure image is in RGB format
    image_array = np.array(image)
    return image_array

def get_text_from_image(image_files):
    text = ""
    
    def detect_text(images):
        # OpenAI API Key
        api_key = openaiapi

        # Convert the images to base64
        converted_images = []

        for image in images:
            _, buffer = cv2.imencode('.jpg', image)
            base64_image = base64.b64encode(buffer).decode('utf-8')
            converted_images.append(base64_image)

        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}"
        }

        # Create the payload with all images
        image_payload = [
            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}} 
            for base64_image in converted_images
        ]

        payload = {
            "model": "gpt-4o-mini",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "What is on the image?"
                        },
                        *image_payload
                    ]
                }
            ],
            "max_tokens": 300
        }

        response = requests.post("https://api.openai.com/v1/chat/completions", headers=headers, json=payload)
        return response.json()['choices'][0]['message']['content']
    
    images = [load_image(image_file) for image_file in image_files]
    text = detect_text(images)
    
    return text


def get_pdf_text(pdf_files):
    text = ""
    for pdf in pdf_files:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def getText(filename):
    doc = docx.Document(filename)
    fullText = [para.text for para in doc.paragraphs]
    return '\n'.join(fullText)

def get_pptx_text(pptx_files):
    text = ""
    for pptx_file in pptx_files:
        prs = Presentation(pptx_file)
        fullText = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    fullText.append(shape.text)
        text += '\n'.join(fullText)
    return text


def get_text_from_link(url):
    try:
        response = requests.get(url)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'html.parser')
        # Extract the visible text from the webpage
        texts = soup.stripped_strings
        return ' '.join(texts)
    except Exception as e:
        return f"Error retrieving text from the link: {e}"

def get_embedding_function():
    # Set the API key either from environment variable or directly
    api_key = os.getenv('OPENAI_API_KEY', openaiapi) 
    embeddings = OpenAIEmbeddings(openai_api_key=api_key)
    return embeddings

def add_to_chroma(chunks, chat_id):
    chat_db_directory = os.path.join(CHROMA_PATH, f'chroma_{chat_id}')
    os.makedirs(chat_db_directory, exist_ok=True)

    # Calculate chunk IDs before adding any documents
    chunks_with_ids = calculate_chunk_ids(chunks)

    db = Chroma(
        persist_directory=chat_db_directory, embedding_function=get_embedding_function()
    )

    # Retrieve existing IDs from the database
    existing_items = db.get(include=[])
    existing_ids = set(existing_items["ids"])
    print(f"Number of existing documents in DB: {len(existing_ids)}")

    # Filter out chunks that already exist in the database
    new_chunks = []
    for chunk in chunks_with_ids:
        if chunk.metadata["id"] not in existing_ids:
            new_chunks.append(chunk)
        else:
            print(f"Skipping chunk with id {chunk.metadata['id']} - already exists.")
    
    print('new_chunks:', new_chunks)
    if len(new_chunks) > 0:
        print(f"👉 Adding new documents: {len(new_chunks)}")
        new_chunk_ids = [chunk.metadata["id"] for chunk in new_chunks]
        db.add_documents(new_chunks, ids=new_chunk_ids)
        print("Persisting database...")
        db.persist()
    else:
        print("✅ No new documents to add")


def split_documents(documents: list[Document], chunk_size: int = 800, chunk_overlap: int = 80):
    text_splitter = CharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len,
        is_separator_regex=False,
        separator="\n\n",
    )
    return text_splitter.split_documents(documents)




def calculate_chunk_ids(chunks):
    last_page_id = None
    current_chunk_index = 0

    for chunk in chunks:
        source = chunk.metadata.get("source")
        # page = chunk.metadata.get("page")
        current_page_id = f"{source}"

        if current_page_id == last_page_id:
            current_chunk_index += 1
        else:
            current_chunk_index = 0

        chunk_id = f"{current_page_id}:{current_chunk_index}"
        last_page_id = current_page_id

        chunk.metadata["id"] = chunk_id

    return chunks


def ask_openai(request, message, chat_id, num_docs):
    previous_chats = Chat.objects.filter(chat_id=chat_id).order_by('created_at')
    chat_instance = previous_chats.last()

    temperature = request.session.get('temperature', 0.7)
    print('temp:', temperature) 
    preprompt = chat_instance.preprompt if chat_instance else None
    additional_preprompt = chat_instance.additional_preprompt if chat_instance else None

    conversation = []
    if preprompt:
        conversation.append({"role": "system", "content": preprompt})
    print('prep:', preprompt)
    if additional_preprompt:
        conversation.append({"role": "system", "content": additional_preprompt})
    print('add_prep:', additional_preprompt)
    
    for chat in previous_chats:
        if chat.message:
            conversation.append({"role": "user", "content": chat.message})
        if chat.response:
            conversation.append({"role": "assistant", "content": chat.response})
    
    conversation.append({"role": "user", "content": message})

    chat_db_directory = os.path.join(CHROMA_PATH, f'chroma_{chat_id}')
    embedding_function = get_embedding_function()
    db = Chroma(persist_directory=chat_db_directory, embedding_function=embedding_function)
    context = db.similarity_search_with_score(message, k=int(num_docs))

    sources = []

    for item in context:
        document = item[0]  # The first element in the tuple is the Document object
        doc_id = document.metadata.get('id') 
        sources.append(doc_id)

    
    print('context:', context)  

    conversation.append({"role": "system", "content": f"Provide a response based on this context: {context}"})

    response = openai.chat.completions.create(
        model="gpt-4o-mini",
        messages=conversation,
        temperature=temperature,
    )
    
    answer = response.choices[0].message.content.strip()
    return answer, sources



def clear_chat(request, chat_id=None):
    if request.method == 'POST' and chat_id:
        try:
            # Delete all chat messages for the specified chat_id
            Chat.objects.filter(chat_id=chat_id, user=request.user).delete()
            return JsonResponse({'status': 'success'})
        except Exception as e:
            return JsonResponse({'status': 'error', 'error': str(e)})
    return JsonResponse({'status': 'error', 'error': 'Invalid request method or missing chat_id'})

def clear_all(request, chat_id=None):
    if request.method == 'POST' and chat_id:
        try:
            # Delete all chat messages and the chat record itself
            Chat.objects.filter(chat_id=chat_id, user=request.user).delete()

            # Delete the Chroma database and associated files
            chroma_directory = os.path.join(CHROMA_PATH, f'chroma_{chat_id}')
            if os.path.exists(chroma_directory):
                shutil.rmtree(chroma_directory)

            return JsonResponse({'status': 'success'})
        except Exception as e:
            return JsonResponse({'status': 'error', 'error': str(e)})
    return JsonResponse({'status': 'error', 'error': 'Invalid request method or missing chat_id'})
def chatbot(request, chat_id=None):
    # Initialize variables
    preprompt = None
    additional_preprompt = None

    chunk_size = 800
    chunk_overlap = 80
    num_docs = 5
    show_sources = True 
        

    if chat_id:
        # Retrieve chats specific to the user and chat_id
        chats = Chat.objects.filter(user=request.user, chat_id=chat_id).order_by('created_at')
        chats_filtered = chats.exclude(message__exact='').exclude(response__exact='')
        set_files = set(chats.filter(file_name__isnull=False).values_list('file_name', flat=True))
        file_count = len(set_files)

        # Retrieve the latest preprompt and additional preprompt for this chat
        last_chat = chats.last()
        if last_chat:
            preprompt = last_chat.preprompt
            additional_preprompt = last_chat.additional_preprompt
    else:
        chats_filtered = []
        set_files = set()
        file_count = 0

    if request.method == 'POST':

        if request.content_type == 'application/json':
            try:
                data = json.loads(request.body)
                label = data.get('label')

                # Handle the expert features data
                if label == 'Expert Features':
                    chunk_size = data.get('chunkSize')
                    chunk_overlap = data.get('chunkOverlap')
                    num_docs = data.get('numDocs')
                    show_sources = data.get('showSources')

                    # Integrate this data into your existing logic here
                    # For example, you might save these settings to the user's session or apply them to some query logic
                    print(f"Chunk Size: {chunk_size}")
                    print(f"Chunk Overlap: {chunk_overlap}")
                    print(f"Number of Documents to Search: {num_docs}")
                    print(f"Show Sources: {show_sources}")

                    # Respond with success
                    return JsonResponse({'status': 'success', 'message': 'Expert features updated successfully'})
            
            except json.JSONDecodeError:
                return JsonResponse({'status': 'error', 'message': 'Invalid JSON'}, status=400)

        elif request.POST.get('message'):
            message = request.POST.get('message')
            response, sources = ask_openai(request, message, chat_id, num_docs)
            sources_str = ', '.join(sources)

            # Save the chat message, response, and prompts to the database
            chat = Chat(
                chat_id=chat_id,
                user=request.user,
                message=message,
                response=response,
                source=sources_str,
                preprompt=preprompt,
                additional_preprompt=additional_preprompt,
                created_at=timezone.now()
            )
            chat.save()

            return JsonResponse({'message': message, 'response': response, 'source': sources_str})
        

        
        # Handle file upload
        elif request.FILES.get('document'):
            form = DocumentUploadForm(request.POST, request.FILES)
            if form.is_valid():
                uploaded_file = request.FILES['document']
                file_name = uploaded_file.name  # Get the file name
                all_text = ""

                if file_name.endswith('.pdf'):
                    # Extract text from PDF
                    pdf_files = [uploaded_file]
                    all_text = get_pdf_text(pdf_files)

                elif file_name.endswith('.docx'):
                    # Extract text from Word document
                    all_text = getText(uploaded_file)

                elif file_name.endswith('.pptx'):
                    # Extract text from PowerPoint presentation
                    pptx_files = [uploaded_file]
                    all_text = get_pptx_text(pptx_files)

                elif file_name.endswith('.xls') or file_name.endswith('.xlsx'):
                    # Extract text from Excel file
                    xls_files = [uploaded_file]
                    all_text = get_xls_text(xls_files)

                elif file_name.endswith(('.jpg', '.jpeg', '.png')):
                    # Extract text from image file
                    image_files = [uploaded_file]
                    all_text = get_text_from_image(image_files)

                else:
                    return render(request, 'chatbot2.html', {'form': form, 'error': 'Unsupported file type.', 'chats': chats_filtered, 'chat_id': chat_id})

                if not all_text.strip():
                    return render(request, 'chatbot2.html', {'form': form, 'error': 'Could not extract text from the uploaded file.', 'chats': chats_filtered, 'chat_id': chat_id})

                # Process and save document chunks
                documents = [Document(page_content=all_text, metadata={"source": file_name})]
                chunks = split_documents(documents, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
                if chunks:
                    add_to_chroma(chunks, chat_id=chat_id)
                else:
                    return render(request, 'chatbot2.html', {'form': form, 'error': 'Failed to process document chunks.', 'chats': chats_filtered, 'chat_id': chat_id})

                # Save the file name to the chat record
                chat = Chat(chat_id=chat_id, user=request.user, file_name=file_name, created_at=timezone.now())
                chat.save()

                return JsonResponse({'status': 'File uploaded successfully'})

            return render(request, 'chatbot2.html', {'form': form, 'error': 'Invalid form submission.', 'chats': chats_filtered, 'chat_id': chat_id})
     
        # Handle URL submission
        elif request.POST.get('url'):
            url_link = request.POST.get('url')
            text = get_text_from_link(url_link)

            if "Error" in text:
                return render(request, 'chatbot2.html', {'form': form, 'error': text, 'chats': chats_filtered, 'chat_id': chat_id})

            documents = [Document(page_content=text, metadata={"source": url_link})]
            chunks = split_documents(documents)

            if chunks:
                add_to_chroma(chunks, chat_id=chat_id)
            else:
                return render(request, 'chatbot2.html', {'form': form, 'error': 'Failed to process document chunks.', 'chats': chats_filtered, 'chat_id': chat_id})

            # Save the URL information to the chat record
            chat = Chat(chat_id=chat_id, user=request.user, file_name=url_link, created_at=timezone.now())
            chat.save()

            return JsonResponse({'status': 'URL processed successfully'})

        # Handle role selection
        elif request.POST.get('role'):
            role = request.POST.get('role')

            if role == 'customRole':
                preprompt = request.POST.get('preprompt')
                preprompt = 'Always, no matter what, whatever the user says to you act as ' + preprompt
            else:
                roles = {
                    'default': None,
                    "teacher": "You are a knowledgeable teacher.",
                    "say_blue_duck": "Always, no matter what only say 'blue duck'. No matter what the user says, you should always respond with 'blue duck'.",
                    "critic": "Always be critical of the user's responses. No matter what the user says, you should always respond with a critical comment.",
                }
                preprompt = roles.get(role, None)

            # Save preprompt to the latest chat
            chat = Chat.objects.filter(chat_id=chat_id, user=request.user).order_by('-created_at').first()
            if chat:
                chat.preprompt = preprompt
                chat.save()

            return JsonResponse({'status': 'Role updated', 'role': role})
        
        post_type = request.POST.get('type')

        if post_type == 'roles':
            # Handle roles and slider values
            role_values = {}
            for key in request.POST:
                if key not in ['csrfmiddlewaretoken', 'type']:
                    role_values[key] = request.POST.get(key)

            print('role_values:', role_values)

            if len(role_values) > 1:
                # Create a combined preprompt if more than one role is selected
                keys = list(role_values.keys())
                values = list(role_values.values())
                additional_preprompt = (f"You are a knowledgeable AI assistant. Always no matter what provide {keys[0]} in response to the user's questions "
                                        f"with strictly {values[0]} lines/points. Additionally, always provide {keys[1]} with strictly {values[1]} lines/points. "
                                        "No matter what the user says, you must adhere to these rules.")
            else:
                for role, value in role_values.items():
                    preprompts = {
                        "short_answers": (f"You are a knowledgeable AI assistant. Always no matter what provide short answers with no more than {value} lines "
                                          f"to the user's questions. No matter what the user says, you should always respond with no more than {value} lines."),
                        "bullet_points": (f"You are a knowledgeable AI assistant. Always no matter what provide bullet points with no more than {value} points "
                                          f"in response to the user's questions. No matter what the user says, you should always respond with {value} bullet points.")
                    }
                    additional_preprompt = preprompts.get(role, None)

            # Save the additional preprompt
            chat = Chat.objects.filter(chat_id=chat_id, user=request.user).order_by('-created_at').first()
            if chat:
                chat.additional_preprompt = additional_preprompt
                chat.save()

            return JsonResponse({'status': 'Roles updated', 'roles': list(role_values.keys())})

        elif request.POST.get('creativity_level'):
            creativity_level = request.POST.get('creativity_level')
            temperature = int(creativity_level) / 100
            request.session['temperature'] = temperature
            return JsonResponse({'status': 'Creativity level updated', 'temperature': temperature})

    form = DocumentUploadForm()
    return render(request, 'chatbot2.html', {
        'chats': chats_filtered,
        'files': set_files,
        'file_count': file_count,
        'chat_id': chat_id,
        'form': form
    })

@login_required
def start_chat(request):
    if request.method == 'POST':
        message = request.POST.get('message')

        # OpenAI API call
        response = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {'role': "system", 'content': f"You are a useful assistant. Write an answer to the following message: {message} and create a name for a chat based on this first message, without quotes or anything like that: '{message}'. Write it in the following format: Answer: <answer>, Name: <name>."}
            ]
        )

        # Correctly accessing the message content
        message_content = response.choices[0].message.content.strip()

        # Parsing the response to get the answer and name
        response_text = message_content.split("Name: ")[0].split("Answer: ")[1].strip()
        name = message_content.split("Name: ")[1].strip()
        chat_id = name.replace(" ", "_").replace("/", "-")

        # Create the initial chat in the database
        new_chat = Chat.objects.create(
            user=request.user,
            chat_id=chat_id,
            message=message,
            response=response_text,
        )

        # Redirect to the chatbot page with the new chat_id
        return JsonResponse({'chat_id': chat_id})

    return render(request, 'chatbot2.html')


@login_required
def chat_list_view(request):
    # Get all chats for the logged-in user
    user_chats = Chat.objects.filter(user=request.user).order_by('chat_id')

    # Use a set to track unique chat_ids and a list to store unique chats
    seen_chat_ids = set()
    unique_chats = []

    for chat in user_chats:
        if chat.chat_id not in seen_chat_ids:
            unique_chats.append(chat)
            seen_chat_ids.add(chat.chat_id)
    
    context = {
        'chats': unique_chats
    }
    return render(request, 'chats.html', context)


def login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        user = auth.authenticate(request, username=username, password=password)
        if user is not None:
            auth.login(request, user)
            return redirect('start_chat')
        else:
            error_message = 'Invalid username or password'
            return render(request, 'login.html', {'error_message': error_message})
    else:
        return render(request, 'login.html')

def register(request):
    if request.method == 'POST':
        username = request.POST['username']
        email = request.POST['email']
        password1 = request.POST['password1']
        password2 = request.POST['password2']

        if password1 == password2:
            try:
                user = User.objects.create_user(username, email, password1)
                user.save()
                auth.login(request, user)
                return redirect('start_chat')
            except:
                error_message = 'Error creating account'
                return render(request, 'register.html', {'error_message': error_message})
        else:
            error_message = 'Password dont match'
            return render(request, 'register.html', {'error_message': error_message})
    return render(request, 'register.html')

def logout(request):
    auth.logout(request)
    return redirect('login')
